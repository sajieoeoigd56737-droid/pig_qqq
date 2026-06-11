from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

OUT_DIR = Path(__file__).parent / "slide_pngs"
OUT_PPTX = Path(__file__).parent / "AP3216C_答辩.pptx"

W_IN, H_IN = 16, 9  # widescreen 16:9

prs = Presentation()
prs.slide_width  = Inches(W_IN)
prs.slide_height = Inches(H_IN)

blank_layout = prs.slide_layouts[6]  # completely blank

pngs = sorted(OUT_DIR.glob("slide_*.png"))
for png in pngs:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(png), 0, 0, Inches(W_IN), Inches(H_IN))
    print(f"  added {png.name}")

prs.save(str(OUT_PPTX))
print(f"\nSaved → {OUT_PPTX}")
